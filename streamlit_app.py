import base64
import json
import os
import tempfile
from io import BytesIO
import subprocess
from pathlib import Path
import shutil
from abc import ABC, abstractmethod
from datetime import date, time

import openai
import pycountry
from gtts import gTTS
from moviepy.editor import (
    AudioFileClip,
    CompositeAudioClip,
    ImageClip,
    concatenate_videoclips,
)
import pdftopng
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from google.cloud import texttospeech
import streamlit as st
import pandas as pd

# --- Constants ---
DEFAULT_SLIDE_DURATION = 3
DEFAULT_OUTPUT_FILENAME = "output_video.mp4"
DEFAULT_GAP_BETWEEN_SLIDES = 1  # seconds

# --- Page Configuration ---
st.set_page_config(
    page_title="Slides2Video",
    page_icon="🎬",
    layout="wide",
    initial_sidebar_state="expanded",
)

# --- Abstract Base Class for TTS Engines ---
class TTS(ABC):
    @abstractmethod
    def synthesize_speech(self, text, temp_dir, slide_index, **kwargs):
        pass

    def display_settings(self, slide_index):
        """Displays UI settings specific to the TTS engine."""
        pass

    @staticmethod
    def get_language_code_from_name(language_name):
        try:
            return pycountry.languages.get(name=language_name).alpha_2 + "-" + pycountry.countries.get(name=language_name.split(" (")[0]).alpha_2
        except KeyError:
            return None

    @staticmethod
    def get_language_name_from_code(language_code):
        try:
            return pycountry.languages.get(alpha_2=language_code.split("-")[0]).name
        except KeyError:
            return "Unknown Language"


# --- TTS Engine Implementations ---
class GTTS_Engine(TTS):
    def synthesize_speech(self, text, temp_dir, slide_index, language="en"):
        tts = gTTS(text=text, lang=language)
        audio_path = os.path.join(temp_dir, f"temp_audio_{slide_index}.mp3")
        tts.save(audio_path)
        return audio_path

    def display_settings(self, slide_index):
        language_options = list(gTTS.tts_langs().keys())
        language_code = st.selectbox("Select Language:", options=language_options, key=f"gtts_language_{slide_index}")
        st.write(f"Language: {gTTS.tts_langs().get(language_code)}")
        return {"language": language_code}  # Return settings as a dictionary


class GoogleCloud_Engine(TTS):
    def __init__(self):
        self.credentials_content = None

    def _set_credentials(self, credentials_content):
        """Sets the Google Cloud credentials."""
        self.credentials_content = credentials_content
        with tempfile.NamedTemporaryFile(mode="w", delete=False) as creds_file:
            creds_file.write(self.credentials_content)
            creds_file.close()
            os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = creds_file.name
            try:
                texttospeech.TextToSpeechClient()
                os.remove(creds_file.name)
            except Exception as e:
                st.warning("Invalid Google Cloud credentials. Please check your input.")
                st.write(e)
                return False  # Indicate failure
            return True  # Indicate success

    def synthesize_speech(self, text, temp_dir, slide_index, voice_name, language_code, speaking_rate, pitch):
        synthesis_input = texttospeech.SynthesisInput(text=text)
        voice = texttospeech.VoiceSelectionParams(name=voice_name, language_code=language_code)
        audio_config = texttospeech.AudioConfig(audio_encoding=texttospeech.AudioEncoding.MP3, speaking_rate=speaking_rate, pitch=pitch)
        client = texttospeech.TextToSpeechClient.from_service_account_json(self.credentials_content)
        response = client.synthesize_speech(input=synthesis_input, voice=voice, audio_config=audio_config)
        audio_path = os.path.join(temp_dir, f"temp_audio_{slide_index}.mp3")

        with open(audio_path, "wb") as out:
            out.write(response.audio_content)
        return audio_path

    def _get_language_options(self):
        """Gets available language codes from Google Cloud TTS."""
        with tempfile.NamedTemporaryFile(mode="w", delete=False) as creds_file:
            creds_file.write(self.credentials_content)
            creds_file.close()
            os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = creds_file.name
            try:
                client = texttospeech.TextToSpeechClient()
                voices = client.list_voices().voices
                language_codes = sorted({voice.language_codes[0] for voice in voices})
                os.remove(creds_file.name)
                return language_codes
            except Exception as e:
                st.error(f"Error fetching languages: {e}")
                return []

    @st.cache_data(show_spinner=False)  # Cache but don't show spinner
    def _get_google_tts_voices(self, language_code="en-US"):
        """Gets available voices for a specific language code."""
        with tempfile.NamedTemporaryFile(mode="w", delete=False) as creds_file:
            creds_file.write(self.credentials_content)
            creds_file.close()
            os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = creds_file.name
            try:
                client = texttospeech.TextToSpeechClient()
                voices = client.list_voices().voices
                voice_options = [
                    (f"{voice.name} ({self.get_language_name_from_code(voice.language_codes[0])}, {texttospeech.SsmlVoiceGender(voice.ssml_gender).name})", voice.name)
                    for voice in voices
                    if voice.language_codes and voice.name and voice.ssml_gender and voice.language_codes[0] == language_code
                ]
                os.remove(creds_file.name)
                return voice_options
            except Exception as e:
                st.error(f"Error initializing Google TTS: {e}")
                return []

    def display_settings(self, slide_index):
        def update_google_credentials():
            if self._set_credentials(st.session_state["google_credentials"]):
                st.session_state["google_credentials_content"] = st.session_state["google_credentials"]
                st.success("Google Cloud credentials validated!")

        google_credentials = st.text_area("Paste Google Cloud Credentials JSON:", key=f"google_credentials", on_change=update_google_credentials)

        # Retrieve credentials from session state if they exist
        if "google_credentials_content" in st.session_state:
            self.credentials_content = st.session_state["google_credentials_content"]
            language_options = self._get_language_options()
            language_code = st.selectbox("Select Language:", options=language_options, key=f"google_language_{slide_index}")
            language_name = self.get_language_name_from_code(language_code)
            st.write(f"Language: {language_name}")

            voice_options = self._get_google_tts_voices(language_code)
            if voice_options:
                selected_voice = st.selectbox("Select a voice:", [voice[0] for voice in voice_options], key=f"google_voice_{slide_index}")
                voice_name = dict(voice_options)[selected_voice]
                st.session_state["selected_voice_name"] = voice_name
                speaking_rate = st.slider("Speaking Rate:", min_value=0.25, max_value=4.0, value=1.0, step=0.25, key=f"speaking_rate_{slide_index}")
                pitch = st.slider("Pitch:", min_value=-20.0, max_value=20.0, value=0.0, step=1.0, key=f"pitch_{slide_index}")

                # Update settings dictionary
                return {"voice_name": voice_name, "language_code": language_code,
                        "speaking_rate": speaking_rate, "pitch": pitch}
            else:
                st.warning("No voices found. Try different credentials or language.")
        return {}


class OpenAI_Engine(TTS):
    def synthesize_speech(self, text, temp_dir, slide_index, voice_name, api_key, model="tts-1-hd"):
        openai.api_key = api_key  # Use API key from settings
        response = openai.audio.speech.create(input=text, voice=voice_name, model=model)
        with tempfile.NamedTemporaryFile(suffix=".mp3", delete=False) as temp_file:
            temp_file.write(response.content)
        return temp_file.name

    def display_settings(self, slide_index):
        def update_openai_api_key():
            api_key = st.session_state["openai_api_key_input"]
            if api_key:
                try:
                    openai.api_key = api_key
                    openai.Engine.list()  
                    st.session_state["openai_api_key"] = api_key
                    st.success("OpenAI API key validated!")

                except openai.error.AuthenticationError:
                    st.warning("Invalid OpenAI API key.")
                except Exception as e:
                    st.warning("An error occurred while validating your API key.")
                    st.write(e)

        openai_api_key = st.text_input("Enter your OpenAI API Key:", key=f"openai_api_key_input", on_change=update_openai_api_key)

        # Only display voice options if the API key is valid
        if "openai_api_key" in st.session_state:
            voice_name = st.selectbox(
                "Select a voice:",
                [
                    "alloy",
                    "echo",
                    "fable",
                    "onyx",
                    "nova",
                    "shimmer",
                ],
                key=f"openai_voice_{slide_index}",
            )
            st.session_state["selected_voice_name"] = voice_name
            model = st.selectbox("Select Model:", ["tts-1", "tts-1-hd"], key=f"openai_model_{slide_index}")
            return {"voice_name": voice_name, "api_key": openai_api_key, "model": model}
        return {}


# --- UI Components ---
class Project:
    def __init__(self, project_name="MyProject"):
        self.project_name = project_name
        self.slides = []  # Initialize with no slides
        self.slide_duration = DEFAULT_SLIDE_DURATION
        self.gap_duration = DEFAULT_GAP_BETWEEN_SLIDES
        self.output_filename = DEFAULT_OUTPUT_FILENAME

    def load_from_file(self, project_data):
        """Loads project data from a dictionary."""
        self.slides = [Slide.from_data(data) for data in project_data.get("slides_data", [])]
        self.slide_duration = project_data.get("slide_duration", DEFAULT_SLIDE_DURATION)
        self.gap_duration = project_data.get("gap_duration", DEFAULT_GAP_BETWEEN_SLIDES)
        self.output_filename = project_data.get("output_filename", DEFAULT_OUTPUT_FILENAME)

    def save_to_file(self):
        """Saves project data to a JSON file."""
        project_data = {
            "slides_data": [slide.to_data() for slide in self.slides],
            "slide_duration": self.slide_duration,
            "gap_duration": self.gap_duration,
            "output_filename": self.output_filename,
            "audio_files": {},
        }

        # Save audio files as base64 encoded strings
        for i, slide in enumerate(self.slides):
            if slide.audio_path and os.path.exists(slide.audio_path):
                with open(slide.audio_path, "rb") as audio_file:
                    audio_base64 = base64.b64encode(audio_file.read()).decode("utf-8")
                    project_data["audio_files"][f"audio_{i}"] = audio_base64

        try:
            with open(f"{self.project_name}.json", "w") as f:
                json.dump(project_data, f)
            st.success(f"Project '{self.project_name}' saved successfully!")
        except Exception as e:
            st.error(f"Error saving project: {e}")


class Slide:
    def __init__(self, pptx_slide=None):
        self.pptx_slide = pptx_slide or self._create_new_slide()
        self.voiceover = ""
        self.audio_path = None
        self.tts_settings = {}

    @classmethod
    def from_data(cls, data):
        """Creates a Slide object from loaded data."""
        slide = cls()
        slide.voiceover = data.get("voiceover", "")
        slide.audio_path = data.get("audio_path", None)
        slide.tts_settings = data.get("tts_settings", {})
        return slide

    def to_data(self):
        """Prepares slide data for saving."""
        return {
            "voiceover": self.voiceover,
            "audio_path": self.audio_path,
            "tts_settings": self.tts_settings,
        }

    def _create_new_slide(self, layout_index=6):  # Default to Blank layout
        """Creates a new slide with the specified layout index."""
        prs = Presentation()
        slide_layout = prs.slide_layouts[layout_index]
        return prs.slides.add_slide(slide_layout)


class Gallery:
    def __init__(self, project):
        self.project = project

    def display(self):
        """Displays the slide gallery and handles new slide creation."""
        cols = st.columns(3)
        for i, slide in enumerate(self.project.slides):
            with cols[i % 3]:
                with st.container():
                    self._display_slide_preview(slide, i)
                    with st.expander(f"Slide {i + 1}", expanded=False):
                        if st.button(f"Edit Slide {i+1}"):
                            st.session_state["editing_slide"] = True
                            st.session_state["edit_slide_index"] = i
                            st.experimental_rerun()

        # --- Add Slide Section ---
        st.header("Add New Slide")
        layout_names = [
            "Title Slide",
            "Title and Content",
            "Section Header",
            "Two Content",
            "Comparison",
            "Title Only",
            "Blank",
            "Content with Caption",
            "Picture with Caption",
        ]
        selected_layout_name = st.selectbox("Select Slide Layout:", layout_names)
        selected_layout_index = layout_names.index(selected_layout_name)
        

        if st.button("Add Slide"):
            self.add_new_slide(selected_layout_index)
            st.experimental_rerun()
    
    def resize_image_to_even_dimensions(self, image):
        width, height = image.size
        even_width = width if width % 2 == 0 else width + 1
        even_height = height if height % 2 == 0 else height + 1
        return image.resize((even_width, even_height))
    
    def _display_slide_preview(self, slide, slide_index):
        """Displays a preview of the slide."""
        prs = Presentation()
        prs.slides.add_slide(slide.pptx_slide)
        
        # Create a temporary PPTX file
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_pptx:
            prs.save(temp_pptx.name)
        
            # Convert the PPTX to PDF using LibreOffice
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
                self.convert_pptx_to_pdf(Path(temp_pptx.name), Path(temp_pdf.name))
        
                # Convert the PDF to PNG
                with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as temp_png:
                    images = pdf2image.convert_from_path(temp_pdf.name, dpi=300)
                    resized_image = self.resize_image_to_even_dimensions(images[0])
                    resized_image.save(temp_png.name, "PNG")
        
                    # Display the PNG preview
                    st.image(Image.open(temp_png.name), width=200, use_column_width=True)
    
        # Play Audio Button - Only show if audio exists
        if slide.audio_path:
            if st.button(f"Play Audio {slide_index + 1}"):
                st.audio(slide.audio_path, format="audio/mp3")

    def add_new_slide(self, layout_index):
        """Adds a new slide with the specified layout index to the project."""
        self.project.slides.append(Slide(pptx_slide=None))  # Create a blank slide
        self.project.slides[-1]._create_new_slide(layout_index=layout_index)

    def slide_to_image(self, slide, slide_index, temp_dir):
        """Saves a slide as an image."""
        image_path = os.path.join(temp_dir, f"slide_{slide_index}.png")
        slide.pptx_slide.save(image_path, format="png")
        return image_path

    def make_chunks(self, size, chunk_size):
        """Splits a sequence into chunks."""
        return [size[i : i + chunk_size] for i in range(0, len(size), chunk_size)]

    def convert_pptx_to_pdf(self, input_pptx: Path):
        """Converts a PPTX file to PDF using LibreOffice."""
        output_pdf_path: Path = Path("downloads") / f"{input_pptx.stem}.pdf"
        try:
            Path("downloads").mkdir(parents=True, exist_ok=True)
            if shutil.which("libreoffice") is None:
                raise FileNotFoundError(
                    "LibreOffice is not installed or not in the PATH. Please install LibreOffice."
                )

            cmd = [
                "libreoffice",
                "--headless",
                "--convert-to",
                "pdf",
                str(input_pptx),
                "--outdir",
                str(output_pdf_path.parent),
            ]
            subprocess.run(
                cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE
            )
            return output_pdf_path

        except Exception as e:
            raise RuntimeError(
                f"Error converting {input_pptx} to PDF: {e}"
            ) from e


class Editor:
    def __init__(self, project, slide_index):
        self.project = project
        self.slide_index = slide_index
        self.slide = self.project.slides[self.slide_index]

    def display(self):
        """Displays the slide editor in a modal."""
        # Modal can be controlled with a session state variable
        if "editing_slide" in st.session_state and st.session_state["editing_slide"]:
            with st.modal(f"Editing Slide {self.slide_index + 1}"):
                self._display_slide_content()
                self._display_voiceover_options()
                if st.button("Done Editing"):
                    st.session_state["editing_slide"] = False
                    st.experimental_rerun()

    def _display_slide_content(self):
        """Displays and allows editing of slide content."""
        st.write(f"**Slide {self.slide_index + 1}**")
        image_stream = BytesIO()
        prs = Presentation()
        prs.slides.add_slide(self.slide.pptx_slide)
        prs.save(image_stream, format="png")
        image_stream.seek(0)
        st.image(Image.open(image_stream), width=600)

        with st.expander(f"Edit Slide {self.slide_index + 1} Content", expanded=True):
            for shape_index, shape in enumerate(self.slide.pptx_slide.shapes):
                if shape.has_text_frame:
                    current_text = shape.text_frame.text
                    new_text = st.text_area(
                        f"Edit text (currently: {current_text}):",
                        value=current_text,
                        key=f"text_{self.slide_index}_{shape_index}",
                    )
                    shape.text_frame.text = new_text
                elif shape.has_table:
                    st.write("Table:")
                    for row_idx, row in enumerate(shape.table.rows):
                        cols = st.columns(len(row.cells))
                        for col_idx, cell in enumerate(row.cells):
                            with cols[col_idx]:
                                current_text = cell.text_frame.text
                                new_text = st.text_input(
                                    f"Cell ({row_idx}, {col_idx})", value=current_text, key=f"table_{self.slide_index}_{shape_index}_{row_idx}_{col_idx}"
                                )
                                cell.text_frame.text = new_text
                elif shape.shape_type == 13:  # Picture
                    st.write(f"Picture {shape_index + 1}")
                    uploaded_image = st.file_uploader(
                        "Replace this picture:", type=["png", "jpg", "jpeg"], key=f"image_upload_{self.slide_index}_{shape_index}"
                    )
                    if uploaded_image is not None:
                        placeholder = prs.slides[0].placeholders[shape_index]
                        placeholder.insert_picture(uploaded_image)
                    st.write("Resize Picture:")
                    col1, col2 = st.columns(2)
                    with col1:
                        width = st.number_input("Width (inches):", value=shape.width.inches, key=f"pic_width_{self.slide_index}_{shape_index}")
                    with col2:
                        height = st.number_input("Height (inches):", value=shape.height.inches, key=f"pic_height_{self.slide_index}_{shape_index}")
                    shape.width = Inches(width)
                    shape.height = Inches(height)

    def _display_voiceover_options(self):
        """Displays and handles voiceover settings and generation."""
        st.write("**Voiceover:**")
        self.slide.voiceover = st.text_area(
            f"Enter/Edit voiceover for Slide {self.slide_index + 1}:",
            key=f"voiceover_{self.slide_index}",
            value=self.slide.voiceover,
        )
        self.slide.tts_settings = self._display_tts_options()
        if self.slide.audio_path:
            st.audio(self.slide.audio_path, format="audio/mp3")

    def _display_tts_options(self):
        """Displays TTS engine selection and settings."""
        st.write("**Voiceover Settings:**")

        # --- Access Selected Engine ---
        selected_tts_engine = None
        if st.session_state["tts_engine"] == "gTTS":
            selected_tts_engine = GTTS_Engine()
        elif st.session_state["tts_engine"] == "Google Cloud":
            selected_tts_engine = GoogleCloud_Engine()
        elif st.session_state["tts_engine"] == "OpenAI":
            selected_tts_engine = OpenAI_Engine()

        # --- Display Engine Settings ---
        if selected_tts_engine:
            settings = selected_tts_engine.display_settings(self.slide_index)

            if st.button(f"Generate Voiceover {self.slide_index + 1}"):
                with st.spinner("Generating voiceover..."):
                    try:
                        if st.session_state["tts_engine"] == "Google Cloud":
                            settings["credentials_content"] = st.session_state.get("google_credentials_content")
                        elif st.session_state["tts_engine"] == "OpenAI":
                            settings["api_key"] = st.session_state.get("openai_api_key")  # Pass API key

                        audio_path = selected_tts_engine.synthesize_speech(
                            text=self.slide.voiceover,
                            temp_dir="temp",
                            slide_index=self.slide_index,
                            **settings
                        )

                        self.slide.audio_path = audio_path
                        st.success("Voiceover generated!")
                        st.audio(audio_path, format="audio/mp3")
                    except Exception as e:
                        st.error(f"Error generating voiceover: {e}")

            return settings

        return {}  # Return empty dictionary if no engine is selected


# --- Cached function to create audio ---
@st.cache_data
def create_audio(text, voice_name, temp_dir, slide_index, tts_engine, credentials_content=None, api_key=None, language="en", speaking_rate=1.0, pitch=0.0, model="tts-1-hd", **kwargs):
    if text:
        if tts_engine == "gTTS":
            audio_path = GTTS_Engine().synthesize_speech(text, temp_dir, slide_index, language=language)
        elif tts_engine == "Google Cloud":
            audio_path = GoogleCloud_Engine().synthesize_speech(text, temp_dir, slide_index, voice_name, language, speaking_rate, pitch, credentials_content)
        elif tts_engine == "OpenAI":
            audio_path = OpenAI_Engine().synthesize_speech(text, temp_dir, slide_index, voice_name, api_key, model)
        else:
            st.error("Invalid TTS engine selected.")
            return None
    else:
        audio_path = os.path.join(temp_dir, f"temp_audio_{slide_index}.mp3")
        try:
            AudioFileClip(Gallery.make_chunks(size=[1], chunk_size=0.1)[0]).write_audiofile(audio_path)
        except Exception as e:
            st.error(f"Error creating silent audio: {e}")
            return None
    return audio_path


# --- Cached Video Generation Function ---
@st.cache_data
def generate_video(slides_data, slide_duration, gap_duration, output_filename, voice_name, credentials_content, tts_engine, end_slide_index=None):
    try:
        slide_duration = int(slide_duration)
        gap_duration = int(gap_duration)
    except ValueError:
        st.error("Slide duration and gap duration must be valid numbers.")
        return

    with tempfile.TemporaryDirectory() as temp_dir:
        return _create_video_from_data(slides_data, slide_duration, gap_duration, output_filename, voice_name, temp_dir, tts_engine, credentials_content, end_slide_index)


def _create_video_from_data(slides_data, slide_duration, gap_duration, output_filename, voice_name, temp_dir, tts_engine, credentials_content, end_slide_index=None):
    image_clips = []
    audio_clips = []
    slides_to_process = slides_data if end_slide_index is None else slides_data[: end_slide_index + 1]
    for i, slide_data in enumerate(slides_to_process):
        slide = slide_data["slide"]
        image_path = Gallery.slide_to_image(slide=slide, slide_index=i, temp_dir=temp_dir)
        image_clips.append(ImageClip(image_path))
        text = slide_data.get("voiceover", "")
        audio_path = create_audio(text, voice_name, temp_dir, i, tts_engine, credentials_content, **slide_data["tts_settings"])
        if audio_path is None: # Handle potential error in create_audio
            return None
        audio_clip = AudioFileClip(audio_path)
        audio_clips.append(audio_clip)
        actual_slide_duration = max(slide_duration, audio_clip.duration)
        image_clips[-1] = image_clips[-1].set_duration(actual_slide_duration)
        if i < len(slides_to_process) - 1:
            image_clips.append(ImageClip(image_path).set_duration(gap_duration))
    combined_audio = CompositeAudioClip(audio_clips)
    final_clips = [img.set_audio(combined_audio) for img in image_clips]
    final_video = concatenate_videoclips(final_clips)
    video_path = os.path.join(temp_dir, output_filename)
    final_video.write_videofile(video_path, fps=24)
    return video_path


# --- Main Streamlit App ---
def main():
    # --- Project Management ---
    if "project" not in st.session_state:
        st.session_state["project"] = Project()
    project = st.session_state["project"]

    # --- Project Settings ---
    project.project_name = st.text_input("Project Name:", value=project.project_name)

    col1, col2 = st.columns(2)
    with col1:
        if st.button("Save Project"):
            project.save_to_file()
    with col2:
        if st.button("Load Project"):
            try:
                with open(f"{project.project_name}.json", "r") as f:
                    project_data = json.load(f)
                project.load_from_file(project_data)
                st.success(f"Project '{project.project_name}' loaded successfully!")
                st.experimental_rerun()  # Rerun to reflect loaded data
            except FileNotFoundError:
                st.error(f"Project '{project.project_name}' not found.")

    # --- TTS Engine Selection ---
    if "tts_engine" not in st.session_state:
        st.session_state["tts_engine"] = "gTTS"

    def update_tts_engine():
        st.session_state["tts_engine"] = st.session_state["tts_engine_selection"]

    tts_engine_name = st.radio(
        "Select TTS Engine:",
        ["gTTS", "Google Cloud", "OpenAI"],
        key="tts_engine_selection",
        horizontal=True,
        on_change=update_tts_engine
    )

    # --- PowerPoint Upload OR New Presentation ---
    uploaded_pptx = st.file_uploader("Upload your PowerPoint presentation (optional)", type=["pptx"])
    create_new = st.checkbox("Create a new presentation", value=True if not uploaded_pptx else False)

    # --- Presentation Logic ---
    if create_new and not uploaded_pptx:
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        project.slides = [Slide(slide) for slide in prs.slides]
    elif uploaded_pptx:
        prs = Presentation(uploaded_pptx)
        project.slides = [Slide(slide) for slide in prs.slides]

    # --- Global Settings ---
    st.header("Global Settings")
    project.slide_duration = st.number_input(
        "Minimum Slide Duration (seconds)", min_value=1, value=project.slide_duration
    )
    project.gap_duration = st.number_input(
        "Gap Between Slides (seconds)", min_value=0, value=project.gap_duration
    )
    project.output_filename = st.text_input("Output Video Filename", value=project.output_filename)

    # --- Slide Gallery ---
    gallery = Gallery(project)
    gallery.display()

    # --- Slide Editor (Modal) ---
    if "editing_slide" in st.session_state and st.session_state["editing_slide"]:
        editor = Editor(project, st.session_state["edit_slide_index"])
        editor.display()

    # --- Delete Slide Button ---
    if len(project.slides) > 1:
        if st.button("Delete Last Slide"):
            project.slides.pop()
            st.experimental_rerun()

    # --- Video Generation ---
    st.header("Create Full Video")
    if st.button("Generate Video"):
        voice_name = st.session_state.get("selected_voice_name")
        credentials_content = st.session_state.get("google_credentials_content")  # For Google Cloud
        api_key = st.session_state.get("openai_api_key")  # For OpenAI
        tts_engine = st.session_state.get("tts_engine", "gTTS")

        with st.spinner("Creating video..."):
            video_path = generate_video(
                [
                    {
                        "slide": slide.pptx_slide,
                        "voiceover": slide.voiceover,
                        "audio_path": slide.audio_path,
                        "tts_settings": slide.tts_settings
                    } for slide in project.slides
                ],
                project.slide_duration,
                project.gap_duration,
                project.output_filename,
                voice_name,
                credentials_content,
                api_key,
                tts_engine,
            )

            if video_path:
                st.success("Video created successfully!")
                st.balloons()
                st.video(video_path)

    # --- Download and Conversion ---
    st.header("Download and Convert")
    col1, col2 = st.columns(2)
    with col1:
        if st.button("Download PPTX"):
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_pptx:
                prs = Presentation()
                for slide in project.slides:
                    prs.slides.add_slide(slide.pptx_slide)
                prs.save(temp_pptx.name)
                st.download_button(
                    "Download PPTX",
                    data=open(temp_pptx.name, "rb").read(),
                    file_name=f"{project.project_name}.pptx",
                    mime="application/pptx",
                )

    with col2:
        if st.button("Convert to PDF"):
            with st.spinner("Converting to PDF..."):
                try:
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_pptx:
                        prs = Presentation()
                        for slide in project.slides:
                            prs.slides.add_slide(slide.pptx_slide)
                        prs.save(temp_pptx.name)
                        pdf_path = gallery.convert_pptx_to_pdf(Path(temp_pptx.name))
                    if pdf_path and Path(pdf_path).exists():
                        st.success("PDF created successfully!")
                        with open(pdf_path, "rb") as pdf_file:
                            st.download_button(
                                label="Download PDF",
                                data=pdf_file,
                                file_name=f"{Path(pdf_path).name}",
                                mime="application/pdf",
                            )
                    else:
                        st.error("PDF conversion failed. Please check the console for errors.")
                except Exception as e:
                    st.error(f"An error occurred during PDF conversion: {e}")


# --- Run the App ---
if __name__ == "__main__":
    main()
