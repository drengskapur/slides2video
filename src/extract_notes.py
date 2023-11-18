#!/usr/bin/env python3
import logging
from pathlib import Path

from pptx import Presentation

logging.basicConfig(level=logging.INFO, format="%(message)s")


def extract_notes(input_pptx):
    try:
        output_dir = Path("assets/notes")
        output_dir.mkdir(parents=True, exist_ok=True)
        presentation = Presentation(str(input_pptx))
        for slide_number, slide in enumerate(presentation.slides, start=1):
            notes_slide = slide.notes_slide
            if notes_slide and notes_slide.notes_text_frame:
                notes_text = notes_slide.notes_text_frame.text.strip()
                output_file = output_dir / f"note_{slide_number}.txt"
                with output_file.open("w", encoding="utf-8") as file:
                    file.write(notes_text)
                logging.info(
                    f"Extracted notes from Slide {slide_number} and saved to {output_file}"
                )
        logging.info(f"Notes extracted from {len(presentation.slides)} slides.")
    except Exception as e:
        logging.info(f"Error extracting notes: {str(e)}")


if __name__ == "__main__":
    input_folder = Path("input")
    pptx_files = list(input_folder.glob("*.pptx"))

    if not pptx_files:
        logging.error("No PowerPoint files found in the 'input' folder.")
    else:
        extract_notes(pptx_files[0])
